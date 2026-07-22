"""
Image Extractor for MarkItDown MCP
Extracts images from various document formats
- With Tesseract installed: Uses OCR for text extraction
- Without Tesseract: Returns compressed grayscale images
"""

import base64
import io
import os
from typing import List, Dict, Optional
from pathlib import Path


def get_tesseract_ocr():
    """Lazy import tesseract_ocr to avoid import errors"""
    try:
        from .tesseract_ocr import tesseract_ocr
        return tesseract_ocr
    except ImportError:
        return None


def compress_large_image(img_bytes: bytes, target_size_mb: float = 0.2) -> bytes:
    """
    Compress image if it's too large

    Args:
        img_bytes: Original image bytes
        target_size_mb: Target size in MB (default: 0.2 = 200KB)

    Returns:
        Compressed image bytes

    Note: Images are converted to grayscale to reduce size
    """
    try:
        from PIL import Image
        import io

        # Load image
        img = Image.open(io.BytesIO(img_bytes))

        # Convert to grayscale (reduces size significantly)
        if img.mode != 'L':
            img = img.convert('L')
            print(f"[Image Extractor] Converted to grayscale")

        # Try different quality levels
        for quality in [85, 75, 65, 55, 45, 35, 25, 15]:
            output = io.BytesIO()
            img.save(output, format='JPEG', quality=quality, optimize=True)
            compressed_size = output.tell()

            # Check if size is acceptable
            if compressed_size <= target_size_mb * 1024 * 1024:
                print(f"[Image Extractor] Compressed: {len(img_bytes) / 1024 / 1024:.1f}MB -> {compressed_size / 1024:.1f}KB (grayscale, quality={quality})")
                output.seek(0)
                return output.read()

        # If still too large, resize
        scale = 0.8
        while scale > 0.3:
            new_width = int(img.width * scale)
            new_height = int(img.height * scale)
            resized = img.resize((new_width, new_height), Image.Resampling.LANCZOS)

            output = io.BytesIO()
            resized.save(output, format='JPEG', quality=65, optimize=True)
            compressed_size = output.tell()

            if compressed_size <= target_size_mb * 1024 * 1024:
                print(f"[Image Extractor] Resized and compressed: {len(img_bytes) / 1024 / 1024:.1f}MB -> {compressed_size / 1024:.1f}KB (grayscale, scale={scale:.1f})")
                output.seek(0)
                return output.read()

            scale -= 0.1

        # Last resort: return highly compressed version
        output = io.BytesIO()
        resized = img.resize((int(img.width * 0.3), int(img.height * 0.3)), Image.Resampling.LANCZOS)
        resized.save(output, format='JPEG', quality=50, optimize=True)
        output.seek(0)
        print(f"[Image Extractor] Heavy compression: {len(img_bytes) / 1024 / 1024:.1f}MB -> {output.tell() / 1024:.1f}KB (grayscale)")
        output.seek(0)
        return output.read()

    except Exception as e:
        print(f"[Image Extractor] Compression failed: {e}")
        return img_bytes


def process_image_with_ocr(img_bytes: bytes, img_info: Dict) -> Optional[Dict]:
    """
    Process image with Tesseract OCR

    Args:
        img_bytes: Raw image bytes
        img_info: Image metadata (page, size, etc.)

    Returns:
        Dict with OCR text or None if Tesseract not available
    """
    tesseract = get_tesseract_ocr()
    if not tesseract or not tesseract.is_available():
        return None

    try:
        # Extract text using Tesseract
        text = tesseract.extract_text_from_image(img_bytes, lang='eng')

        if text is not None:
            return {
                **img_info,
                'ocr_text': text,
                'method': 'tesseract_ocr',
                'has_text': len(text) > 0
            }

    except Exception as e:
        print(f"[Image Extractor] OCR failed: {e}")

    return None


def extract_images_from_pdf(file_path: str) -> List[Dict[str, any]]:
    """Extract images from PDF file"""
    try:
        import pdfplumber
        from PIL import Image

        images = []

        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_images = page.images

                for img_idx, img_info in enumerate(page_images):
                    try:
                        # Try to extract image using bbox
                        x0 = img_info.get('x0', 0)
                        y0 = img_info.get('top', 0)
                        x1 = img_info.get('x1', x0 + img_info.get('width', 100))
                        y1 = img_info.get('bottom', y0 + img_info.get('height', 100))

                        # Skip invalid dimensions
                        if x1 <= x0 or y1 <= y0:
                            continue

                        # Crop and render the region
                        bbox = (x0, y0, x1, y1)
                        cropped = page.within_bbox(bbox)
                        img = cropped.to_image(resolution=150)

                        # Convert to PIL Image
                        img_bytes = io.BytesIO()
                        img.original.save(img_bytes, format='PNG')
                        img_bytes.seek(0)

                        original_bytes = img_bytes.read()
                        original_size = len(original_bytes)

                        base_info = {
                            'page': page_num,
                            'index': img_idx,
                            'width': int(x1 - x0),
                            'height': int(y1 - y0),
                            'original_size': original_size
                        }

                        # Try OCR first
                        ocr_result = process_image_with_ocr(original_bytes, base_info)
                        if ocr_result:
                            images.append(ocr_result)
                            continue

                        # Fallback: Compress if larger than 200KB
                        if original_size > 200 * 1024:
                            compressed_bytes = compress_large_image(original_bytes, target_size_mb=0.2)
                            base64_img = base64.b64encode(compressed_bytes).decode('utf-8')
                            img_format = 'jpeg'  # Compressed images are JPEG
                        else:
                            base64_img = base64.b64encode(original_bytes).decode('utf-8')
                            img_format = 'png'

                        images.append({
                            **base_info,
                            'base64': base64_img,
                            'format': img_format,
                            'compressed': original_size > 200 * 1024,
                            'method': 'grayscale_compression'
                        })

                    except Exception as e:
                        print(f"[Image Extractor] Warning: Failed to extract image from page {page_num}: {e}")
                        continue

        return images

    except ImportError:
        print("[Image Extractor] Error: pdfplumber not installed")
        return []
    except Exception as e:
        print(f"[Image Extractor] Error extracting PDF images: {e}")
        return []


def extract_images_from_docx(file_path: str) -> List[Dict[str, any]]:
    """Extract images from DOCX file"""
    try:
        from docx import Document

        images = []
        doc = Document(file_path)

        # Extract images from document relationships
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    img_data = rel.target_part.blob

                    # Determine image format from content type
                    content_type = rel.target_part.content_type
                    img_format = content_type.split('/')[-1] if '/' in content_type else 'png'

                    base_info = {
                        'index': len(images),
                        'format': img_format,
                        'original_size': len(img_data)
                    }

                    # Try OCR first
                    ocr_result = process_image_with_ocr(img_data, base_info)
                    if ocr_result:
                        images.append(ocr_result)
                        continue

                    # Fallback: Return base64
                    base64_img = base64.b64encode(img_data).decode('utf-8')
                    images.append({
                        **base_info,
                        'base64': base64_img,
                        'method': 'grayscale_compression'
                    })

                except Exception as e:
                    print(f"[Image Extractor] Warning: Failed to extract DOCX image: {e}")
                    continue

        return images

    except ImportError:
        print("[Image Extractor] Error: python-docx not installed")
        return []
    except Exception as e:
        print(f"[Image Extractor] Error extracting DOCX images: {e}")
        return []


def extract_images_from_pptx(file_path: str) -> List[Dict[str, any]]:
    """Extract images from PPTX file"""
    try:
        from pptx import Presentation

        images = []
        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    try:
                        img_data = shape.image.blob
                        img_format = shape.image.ext.lstrip('.')

                        base_info = {
                            'slide': slide_num,
                            'index': len(images),
                            'format': img_format,
                            'original_size': len(img_data)
                        }

                        # Try OCR first
                        ocr_result = process_image_with_ocr(img_data, base_info)
                        if ocr_result:
                            images.append(ocr_result)
                            continue

                        # Fallback: Return base64
                        base64_img = base64.b64encode(img_data).decode('utf-8')
                        images.append({
                            **base_info,
                            'base64': base64_img,
                            'method': 'grayscale_compression'
                        })

                    except Exception as e:
                        print(f"[Image Extractor] Warning: Failed to extract PPTX image from slide {slide_num}: {e}")
                        continue

        return images

    except ImportError:
        print("[Image Extractor] Error: python-pptx not installed")
        return []
    except Exception as e:
        print(f"[Image Extractor] Error extracting PPTX images: {e}")
        return []


def extract_images_from_xlsx(file_path: str) -> List[Dict[str, any]]:
    """Extract images from XLSX file"""
    try:
        from openpyxl import load_workbook

        images = []
        wb = load_workbook(file_path)

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]

            if hasattr(sheet, '_images') and sheet._images:
                for img_idx, img in enumerate(sheet._images):
                    try:
                        img_data = img._data()

                        base_info = {
                            'sheet': sheet_name,
                            'index': img_idx,
                            'format': 'png',
                            'original_size': len(img_data)
                        }

                        # Try OCR first
                        ocr_result = process_image_with_ocr(img_data, base_info)
                        if ocr_result:
                            images.append(ocr_result)
                            continue

                        # Fallback: Return base64
                        base64_img = base64.b64encode(img_data).decode('utf-8')
                        images.append({
                            **base_info,
                            'base64': base64_img,
                            'method': 'grayscale_compression'
                        })

                    except Exception as e:
                        print(f"[Image Extractor] Warning: Failed to extract XLSX image from {sheet_name}: {e}")
                        continue

        return images

    except ImportError:
        print("[Image Extractor] Error: openpyxl not installed")
        return []
    except Exception as e:
        print(f"[Image Extractor] Error extracting XLSX images: {e}")
        return []


def extract_images_from_html(file_path: str) -> List[Dict[str, any]]:
    """Extract images from HTML file"""
    try:
        from bs4 import BeautifulSoup

        images = []

        with open(file_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')

        base_url = Path(file_path).parent.as_uri()

        for img_idx, img_tag in enumerate(soup.find_all('img')):
            src = img_tag.get('src')
            if not src:
                continue

            try:
                img_data = None

                # Handle data URIs
                if src.startswith('data:image'):
                    if 'base64,' in src:
                        base64_data = src.split('base64,')[1]
                        img_data = base64.b64decode(base64_data)

                # Handle local files
                else:
                    from urllib.parse import urljoin
                    img_url = urljoin(base_url, src)

                    if img_url.startswith('file://'):
                        img_path = img_url.replace('file://', '')
                        with open(img_path, 'rb') as img_file:
                            img_data = img_file.read()

                if not img_data:
                    continue

                base_info = {
                    'index': img_idx,
                    'format': Path(src).suffix.lstrip('.') or 'png',
                    'original_size': len(img_data)
                }

                # Try OCR first
                ocr_result = process_image_with_ocr(img_data, base_info)
                if ocr_result:
                    images.append(ocr_result)
                    continue

                # Fallback: Return base64
                base64_img = base64.b64encode(img_data).decode('utf-8')
                images.append({
                    **base_info,
                    'base64': base64_img,
                    'method': 'grayscale_compression'
                })

            except Exception as e:
                print(f"[Image Extractor] Warning: Failed to extract HTML image {src}: {e}")
                continue

        return images

    except ImportError:
        print("[Image Extractor] Error: beautifulsoup4 not installed")
        return []
    except Exception as e:
        print(f"[Image Extractor] Error extracting HTML images: {e}")
        return []


def extract_image_from_file(file_path: str) -> List[Dict[str, any]]:
    """Extract image from image file (PNG, JPG, etc.)"""
    try:
        with open(file_path, 'rb') as f:
            img_data = f.read()

        base_info = {
            'index': 0,
            'format': Path(file_path).suffix.lstrip('.'),
            'original_size': len(img_data)
        }

        # Try OCR first
        ocr_result = process_image_with_ocr(img_data, base_info)
        if ocr_result:
            return [ocr_result]

        # Fallback: Return base64
        base64_img = base64.b64encode(img_data).decode('utf-8')
        return [{
            **base_info,
            'base64': base64_img,
            'method': 'grayscale_compression'
        }]

    except Exception as e:
        print(f"[Image Extractor] Error reading image file: {e}")
        return []


def extract_images(file_path: str) -> List[Dict[str, any]]:
    """
    Extract images from various document formats

    Supported formats:
    - PDF (.pdf) - with Tesseract OCR or grayscale compression
    - DOCX (.docx)
    - PPTX (.pptx)
    - XLSX (.xlsx)
    - HTML (.html, .htm)
    - Images (.png, .jpg, .jpeg, .gif, .bmp)

    Returns:
        List of dicts with keys:
        - Tesseract mode: ocr_text, method='tesseract_ocr'
        - Compression mode: base64, format, compressed, method='grayscale_compression'
    """
    file_ext = Path(file_path).suffix.lower()

    if file_ext == '.pdf':
        return extract_images_from_pdf(file_path)
    elif file_ext == '.docx':
        return extract_images_from_docx(file_path)
    elif file_ext == '.pptx':
        return extract_images_from_pptx(file_path)
    elif file_ext in ['.xlsx', '.xls']:
        return extract_images_from_xlsx(file_path)
    elif file_ext in ['.html', '.htm']:
        return extract_images_from_html(file_path)
    elif file_ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp']:
        return extract_image_from_file(file_path)
    else:
        print(f"[Image Extractor] Unsupported format: {file_ext}")
        return []


def format_images_as_markdown(images: List[Dict[str, any]], max_images: int = None) -> str:
    """
    Format extracted images as Markdown

    Args:
        images: List of image dicts from extract_images()
        max_images: Maximum number of images to include (None = all)

    Returns:
        Markdown string with embedded images or OCR text
    """
    if not images:
        return ""

    if max_images:
        images = images[:max_images]

    markdown_parts = []

    for i, img in enumerate(images, 1):
        method = img.get('method', 'unknown')

        # Add position info if available
        position_info = []
        if 'page' in img:
            position_info.append(f"Page {img['page']}")
        if 'slide' in img:
            position_info.append(f"Slide {img['slide']}")
        if 'sheet' in img:
            position_info.append(f"Sheet: {img['sheet']}")

        position_str = f" ({', '.join(position_info)})" if position_info else ""

        if method == 'tesseract_ocr':
            # OCR text output
            ocr_text = img.get('ocr_text', '')
            if ocr_text:
                markdown_parts.append(
                    f"\n### Image {i}{position_str} - OCR Text\n\n```\n{ocr_text}\n```\n"
                )
            else:
                markdown_parts.append(
                    f"\n### Image {i}{position_str}\n\n*No text detected*\n"
                )
        else:
            # Base64 image output
            img_format = img.get('format', 'png')
            base64_data = img.get('base64', '')
            markdown_parts.append(
                f"\n![Image {i}{position_str}](data:image/{img_format};base64,{base64_data})\n"
            )

    return "\n".join(markdown_parts)
